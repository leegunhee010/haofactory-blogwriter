# -*- coding: utf-8 -*-
"""카드뉴스 PIL 엔진 — 브랜드별 템플릿을 순수 Pillow로 재현. PowerPoint 불필요.
브랜드 에셋: brands/<id>/cards/ (layout.json + s*.png), 폰트: assets/fonts 공용.
extract_template()으로 어떤 PPTX든 에셋+layout으로 추출(브랜드 등록용)."""
import io, os, re, json, random
from PIL import Image, ImageOps, ImageEnhance, ImageFilter, ImageDraw, ImageFont
Image.MAX_IMAGE_PIXELS = None   # 사용자 본인 사진 — 대형 이미지 허용(압축폭탄 방어 해제)
try:
    import numpy as _np
except Exception:
    _np = None

HERE = os.path.dirname(os.path.abspath(__file__))
FONTDIR = os.path.join(HERE, "assets", "fonts")
_HEAD_FONT = os.path.join(FONTDIR, "Pretendard-Bold.otf")
_LABEL_FONT = os.path.join(FONTDIR, "GmarketSansTTFBold.ttf")
_BODY_FONT = os.path.join(FONTDIR, "Pretendard-Regular.otf")
if not os.path.isfile(_BODY_FONT):
    _BODY_FONT = _HEAD_FONT

# 템플릿 원본 폰트명 → 동봉 폰트파일 (원본 디자인 그대로 재현). 없으면 기본폰트 fallback.
FONT_FILES = {
    "gmarket sans bold": "GmarketSansTTFBold.ttf", "gmarket sans": "GmarketSansTTFBold.ttf",
    "gmarket sans medium": "GmarketSansTTFMedium.ttf", "gmarket sans ttf medium": "GmarketSansTTFMedium.ttf",
    "210 sanullim regular": "210Sanullim.ttf", "210 sanullim": "210Sanullim.ttf",
    "pretendard bold": "Pretendard-Bold.otf", "pretendard": "Pretendard-Bold.otf",
    "pretendard medium": "Pretendard-Medium.otf", "pretendard semibold": "Pretendard-SemiBold.otf",
    "pretendard regular": "Pretendard-Regular.otf", "pretendard black": "Pretendard-Black.otf",
    "nexon lv1 gothic otf bold": "NEXONLv1GothicOTFBold.otf", "nexon lv1 gothic": "NEXONLv1GothicOTFBold.otf",
    "gangwonedupower": "GangwonEduPowerExtraBoldA.otf", "gangwonedu": "GangwonEduPowerExtraBoldA.otf",
    "paperlogy 7 bold": "Paperlogy-7Bold.ttf", "paperlogy": "Paperlogy-7Bold.ttf",
}


def _font_for(name, default_path):
    """원본 폰트명으로 동봉 폰트파일 찾기. 없으면 default_path."""
    if not name:
        return default_path
    key = re.sub(r"\s+", " ", str(name).strip().lower())
    fn = FONT_FILES.get(key)
    if not fn:  # 느슨한 매칭(앞부분 일치)
        for k, v in FONT_FILES.items():
            if key.startswith(k) or k.startswith(key):
                fn = v; break
    if fn:
        p = os.path.join(FONTDIR, fn)
        if os.path.isfile(p):
            return p
    # assets/fonts 에서 파일명으로 매칭(PPTX에서 추출한 임베드 폰트)
    try:
        for f in os.listdir(FONTDIR):
            if not f.lower().endswith((".ttf", ".otf")):
                continue
            stem = re.sub(r"\s+", " ", os.path.splitext(f)[0].strip().lower())
            if stem == key:
                return os.path.join(FONTDIR, f)
    except Exception:
        pass
    return default_path


def embed_fonts_from_pptx(pptx_path):
    """PPTX에 임베드된 폰트를 assets/fonts 로 추출(타이페이스명으로 저장). 추출된 폰트명 리스트 반환."""
    import zipfile
    from pptx.oxml.ns import qn as _qn
    out = []
    z = zipfile.ZipFile(pptx_path)
    try:
        pres = z.read("ppt/presentation.xml")
        rels = z.read("ppt/_rels/presentation.xml.rels")
    except KeyError:
        return out
    import xml.etree.ElementTree as ET
    P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    rid2tgt = {}
    for rel in ET.fromstring(rels):
        rid2tgt[rel.get("Id")] = rel.get("Target")
    for ef in ET.fromstring(pres).iter(P + "embeddedFont"):
        fontEl = ef.find(P + "font")
        if fontEl is None:
            continue
        typeface = fontEl.get("typeface")
        rid = None
        for tag in ("regular", "bold", "italic", "boldItalic"):
            el = ef.find(P + tag)
            if el is not None and el.get(R + "id"):
                rid = el.get(R + "id"); break
        if not (typeface and rid and rid in rid2tgt):
            continue
        import struct
        tgt = rid2tgt[rid]
        arc = "ppt/" + tgt[3:] if tgt.startswith("../") else "ppt/" + tgt
        try:
            data = z.read(arc.replace("ppt/ppt/", "ppt/"))
        except KeyError:
            continue
        sigs = (b"OTTO", bytes([0, 1, 0, 0]), b"true", b"ttcf")
        real = data
        if data[:4] not in sigs and len(data) >= 16:   # PowerPoint 임베드 = EOT → 헤더 벗기기
            fds = struct.unpack("<I", data[4:8])[0]
            flags = struct.unpack("<I", data[12:16])[0]
            if 0 < fds <= len(data) and not (flags & 0x4):   # 비압축 EOT
                cand = data[len(data) - fds:]
                if cand[:4] in sigs:
                    real = cand
        if real[:4] not in sigs:
            continue                                   # 못 읽는 형식(압축 EOT 등) 건너뜀
        ext = ".otf" if real[:4] == b"OTTO" else ".ttf"
        safe = re.sub(r'[\\/:*?"<>|]', "", typeface).strip()
        dst = os.path.join(FONTDIR, safe + ext)
        with open(dst, "wb") as f:
            f.write(real)
        out.append(typeface)
    return out

PALETTE = {"orange": (253, 111, 34), "yellow": (242, 178, 52), "green": (80, 178, 40),
           "blue": (108, 121, 214), "pink": (224, 106, 140)}

_layout_cache = {}
_recolor_cache = {}
_font_cache = {}


def _layout(assets_dir):
    if assets_dir not in _layout_cache:
        _layout_cache[assets_dir] = json.load(open(os.path.join(assets_dir, "layout.json"), encoding="utf-8"))
    return _layout_cache[assets_dir]


def _font(path, px):
    k = (path, px)
    if k not in _font_cache:
        try:
            _font_cache[k] = ImageFont.truetype(path, px)
        except Exception:
            _font_cache[k] = ImageFont.truetype(_HEAD_FONT, px)   # 못 읽는 폰트면 기본폰트
    return _font_cache[k]


def _remap(im, theme):
    """주황 장식 → 테마색 (흰→테마 tint 보존)."""
    im = im.convert("RGBA")
    if _np is not None:
        a = _np.asarray(im).astype("float32")
        mn = a[..., :3].min(axis=2)
        t = ((255.0 - mn) / (255 - 34)).clip(0, 1)[..., None]
        rgb = 255.0 * (1 - t) + _np.array(theme, "float32") * t
        out = _np.dstack([rgb, a[..., 3:4]]).clip(0, 255).astype("uint8")
        return Image.fromarray(out, "RGBA")
    px = im.load(); W, H = im.size
    for y in range(H):
        for x in range(W):
            r, g, b, al = px[x, y]
            if al < 5:
                continue
            tt = max(0.0, min(1.0, (255 - min(r, g, b)) / (255 - 34)))
            px[x, y] = (round(255 * (1 - tt) + theme[0] * tt), round(255 * (1 - tt) + theme[1] * tt),
                        round(255 * (1 - tt) + theme[2] * tt), al)
    return im


def _asset(assets_dir, name, box, theme, recolor):
    key = (assets_dir, name, theme if recolor else None)
    if key not in _recolor_cache:
        im = Image.open(os.path.join(assets_dir, name)).convert("RGBA")
        if recolor:
            im = _remap(im, theme)
        _recolor_cache[key] = im
    im = _recolor_cache[key]
    if im.size != (box[2], box[3]):
        im = im.resize((max(1, box[2]), max(1, box[3])), Image.LANCZOS)
    return im


def _paste(canvas, im, x, y):
    if im.mode != "RGBA":
        im = im.convert("RGBA")
    layer = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    sx = max(0, -x); sy = max(0, -y)
    if sx or sy:
        im = im.crop((sx, sy, im.width, im.height)); x += sx; y += sy
    layer.paste(im, (x, y))
    return Image.alpha_composite(canvas, layer)


def _fit(path, W, H, seed, center=None, zoom=1.0):
    rnd = random.Random(seed)
    im = Image.open(path)
    if im.mode in ("RGBA", "LA", "P"):
        im = im.convert("RGBA"); bg = Image.new("RGBA", im.size, (255, 255, 255, 255))
        im = Image.alpha_composite(bg, im).convert("RGB")
    else:
        im = im.convert("RGB")
    im = ImageOps.exif_transpose(im)
    for E in (ImageEnhance.Brightness, ImageEnhance.Contrast, ImageEnhance.Color):
        im = E(im).enhance(rnd.uniform(0.96, 1.05))
    im = im.filter(ImageFilter.UnsharpMask(1.2, rnd.randint(40, 70), 2))
    cx = 0.5 if center is None else max(0.0, min(1.0, center[0]))
    cy = rnd.uniform(0.42, 0.58) if center is None else max(0.0, min(1.0, center[1]))
    zoom = max(1.0, min(2.5, float(zoom or 1.0)))
    if zoom > 1.0:
        w, h = im.size; cw, ch = int(w / zoom), int(h / zoom)
        left = int((w - cw) * cx); top = int((h - ch) * cy)
        im = im.crop((left, top, left + cw, top + ch))
    return ImageOps.fit(im, (W, H), Image.LANCZOS, centering=(cx, cy))


def _round(im, rad):
    rad = max(0, int(rad))
    if rad == 0:
        return im.convert("RGBA")
    im = im.convert("RGBA"); W, H = im.size
    m = Image.new("L", (W, H), 0)
    ImageDraw.Draw(m).rounded_rectangle([0, 0, W - 1, H - 1], radius=rad, fill=255)
    im.putalpha(m)
    return im


def _circle(im):
    im = im.convert("RGBA"); W, H = im.size
    m = Image.new("L", (W, H), 0)
    ImageDraw.Draw(m).ellipse([0, 0, W - 1, H - 1], fill=255)
    im.putalpha(m)
    return im


def _accent_from_image(path, light_bg=True):
    """사진에서 가장 또렷한 대표색 추출(글자색용). 밝은 배경 위 가독 위해 명도 캡."""
    import colorsys
    try:
        im = Image.open(path).convert("RGB")
    except Exception:
        return None
    im = ImageOps.exif_transpose(im).resize((48, 48))
    scored = []
    for r, g, b in im.getdata():
        h, s, v = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
        scored.append((s * (1 - abs(v - 0.55)), (r, g, b)))   # 채도 높고 중간밝기 우선
    scored.sort(key=lambda x: x[0], reverse=True)
    top = [c for _, c in scored[:160]]
    if not top:
        return None
    r = sum(c[0] for c in top) // len(top)
    g = sum(c[1] for c in top) // len(top)
    b = sum(c[2] for c in top) // len(top)
    h, s, v = colorsys.rgb_to_hsv(r / 255, g / 255, b / 255)
    if light_bg:
        s = min(1.0, s * 1.25); v = min(v, 0.6)   # 밝은배경 가독: 채도↑·명도 캡
    r, g, b = [int(round(x * 255)) for x in colorsys.hsv_to_rgb(h, s, v)]
    return (r, g, b)


def _random_accent():
    """생성마다 다른 선명한 색(밝은배경 위 가독·장식 recolor 겸용 중간명도)."""
    import colorsys
    h = random.random()
    s = random.uniform(0.45, 0.78)
    v = random.uniform(0.52, 0.70)
    r, g, b = colorsys.hsv_to_rgb(h, s, v)
    return (int(round(r * 255)), int(round(g * 255)), int(round(b * 255)))


def _shorten_title(t, maxlen=24):
    """제목이 너무 길면 어절·구두점 경계에서 줄임(표지 글자수 한계)."""
    t = " ".join(str(t).split())
    if len(t) <= maxlen:
        return t
    cut = t[:maxlen]
    best = -1
    for sep in (" ", ",", "·", "—", "-"):
        i = cut.rfind(sep)
        if i > best and i > maxlen * 0.55:
            best = i
    if best > 0:
        cut = cut[:best]
    return cut.strip().rstrip(",·-— ")


def _draw_text(canvas, box, text, font_path, size_pt, theme, align, pt2px, oneline=False):
    if not text:
        return canvas
    x, y, w, h = box
    col = (theme[0], theme[1], theme[2], 255)
    px = max(10, int(round(size_pt * pt2px)))
    layer = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)

    def build(fnt):   # oneline=줄바꿈 금지(폭 맞게 축소만). 아니면 \n 분리 후 폭 맞춰 단어 줄바꿈
        if oneline:
            return [" ".join(str(text).split())]
        out = []
        for raw in str(text).split("\n"):
            raw = raw.strip()
            if not raw:
                out.append(""); continue
            if d.textlength(raw, font=fnt) <= w or " " not in raw:
                out.append(raw); continue
            cur = ""
            for wd in raw.split(" "):
                tt = (cur + " " + wd).strip()
                if d.textlength(tt, font=fnt) <= w:
                    cur = tt
                else:
                    if cur:
                        out.append(cur)
                    cur = wd
            if cur:
                out.append(cur)
        return out
    fnt = _font(font_path, px)
    lines = build(fnt)
    while px > 12 and lines and max((d.textlength(ln, font=fnt) for ln in lines if ln), default=0) > w:
        px -= 2; fnt = _font(font_path, px); lines = build(fnt)
    asc, desc = fnt.getmetrics(); lh = asc + desc
    cy = y + (h - lh * len(lines)) / 2
    for ln in lines:
        tw = d.textlength(ln, font=fnt) if ln else 0
        tx = x + w - tw if align == "right" else (x + (w - tw) / 2 if align == "center" else x)
        d.text((tx, cy), ln, font=fnt, fill=col)
        cy += lh
    return Image.alpha_composite(canvas, layer)


def _two_lines(text):
    """텍스트를 길이 균형 맞춰 2줄로 강제 분할(공백 기준)."""
    text = " ".join(str(text).split())
    words = text.split(" ")
    if len(words) <= 1:
        return text
    best_i, best_diff = 1, 10 ** 9
    for i in range(1, len(words)):
        a = " ".join(words[:i]); b = " ".join(words[i:])
        diff = abs(len(a) - len(b))
        if diff < best_diff:
            best_diff, best_i = diff, i
    return " ".join(words[:best_i]) + "\n" + " ".join(words[best_i:])


def _draw_kw_title(canvas, box, keyword, title, font_path, kw_size, title_size, color, pt2px, align="right", gap=0, dy=0):
    """키워드(큰글씨) + 제목(작은글씨)을 각자 원본 사이즈로, 정렬·세로가운데 그대로 렌더."""
    x, y, w, h = box
    col = (color[0], color[1], color[2], 255)
    layer = Image.new("RGBA", canvas.size, (0, 0, 0, 0)); d = ImageDraw.Draw(layer)
    kf = _font(font_path, max(10, int(round(kw_size * pt2px))))
    tf = _font(font_path, max(10, int(round(title_size * pt2px))))

    def wrap(text, fnt):
        out = []
        for raw in str(text).split("\n"):
            raw = raw.strip()
            if not raw:
                continue
            if d.textlength(raw, font=fnt) <= w or " " not in raw:
                out.append(raw); continue
            cur = ""
            for wd in raw.split(" "):
                t = (cur + " " + wd).strip()
                if d.textlength(t, font=fnt) <= w:
                    cur = t
                else:
                    if cur:
                        out.append(cur)
                    cur = wd
            if cur:
                out.append(cur)
        return out
    klines = wrap(keyword, kf) if keyword else []
    # 제목은 무조건 한 줄. 폭을 넘으면 끝 어절부터 떼어 축약, 그래도 길면 글자 단위로 컷
    tlines = []
    if title:
        one = " ".join(str(title).split())
        while d.textlength(one, font=tf) > w and " " in one:
            one = one.rsplit(" ", 1)[0]
        while d.textlength(one, font=tf) > w and len(one) > 1:
            one = one[:-1]
        one = one.strip().rstrip(",·-— ")
        if one:
            tlines = [one]
    ka, kd = kf.getmetrics(); klh = ka + kd
    ta, td = tf.getmetrics(); tlh = ta + td
    g = gap if (klines and tlines) else 0
    total = klh * len(klines) + tlh * len(tlines) + g
    cy = y + (h - total) / 2 + dy                       # dy = 전체 아래로 이동
    def drawline(ln, fnt):
        tw = d.textlength(ln, font=fnt)
        tx = x + w - tw if align == "right" else (x + (w - tw) / 2 if align == "center" else x)
        d.text((tx, cy), ln, font=fnt, fill=col)
    for ln in klines:
        drawline(ln, kf); cy += klh
    cy += g                                             # gap = 키워드와 제목 줄 사이 간격
    for ln in tlines:
        drawline(ln, tf); cy += tlh
    return Image.alpha_composite(canvas, layer)


def render_card(slide_idx, photo_path, headline, theme, out_path, assets_dir,
                center=None, zoom=1.0, seed="", subtitle="", body="", title="", accent=None):
    lay = _layout(assets_dir)
    SIZE = lay["size"]
    pt2px = SIZE / float(lay.get("pt_height", 810))
    theme_recolor = lay.get("theme_recolor", True)   # 기존(하오팩토리) layout은 키 없음 → True
    auto_theme = lay.get("auto_theme", False)         # 브랜드색 전체를 이미지색으로 치환
    pround = lay.get("photo_round", 0.045)
    items = lay["slides"][slide_idx]
    # 대표색: accent 전달되면 통일색 사용, 없으면 이 카드 사진에서 추출
    acc = None
    if auto_theme or any(it.get("auto_color") for it in items):
        if accent:
            acc = tuple(accent)
        elif photo_path and os.path.isfile(photo_path):
            acc = _accent_from_image(photo_path)
    eff_theme = acc if (auto_theme and acc) else theme   # 장식 recolor용 테마
    sb = lay.get("slide_bg")
    bg = (sb[slide_idx] if sb and slide_idx < len(sb) else lay.get("bg", "FFFFFF")) or "FFFFFF"
    if auto_theme and acc:
        br, bgc, bb = (int(bg[i:i+2], 16) for i in (0, 2, 4))
        if max(br, bgc, bb) - min(br, bgc, bb) > 30:     # 채도있는(브랜드색) 배경만 치환, 흰/검 유지
            bg = "%02X%02X%02X" % acc
    canvas = Image.new("RGBA", (SIZE, SIZE), tuple(int(bg[i:i+2], 16) for i in (0, 2, 4)) + (255,))

    def textcol(it):
        if it.get("auto_color") and acc:
            return acc
        if not theme_recolor and it.get("color"):
            c = it["color"]
            return (int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16))
        return eff_theme
    for it in items:
        role = it["role"]; box = it["box"]
        if role == "photo":
            if photo_path and os.path.isfile(photo_path):
                if it.get("shape") == "circle":
                    d = min(box[2], box[3])
                    ph = _fit(photo_path, d, d, seed or "0", center=center, zoom=zoom)
                    ph = _circle(ph)
                    canvas = _paste(canvas, ph, box[0] + (box[2] - d) // 2, box[1] + (box[3] - d) // 2)
                else:
                    ph = _fit(photo_path, box[2], box[3], seed or "0", center=center, zoom=zoom)
                    bl = it.get("blur")
                    if bl:
                        ph = ph.filter(ImageFilter.GaussianBlur(bl if isinstance(bl, (int, float)) and bl > 1 else 28))
                    ph = _round(ph, int(min(box[2], box[3]) * pround))
                    canvas = _paste(canvas, ph, box[0], box[1])
        elif role in ("deco", "line", "logo"):
            canvas = _paste(canvas, _asset(assets_dir, it["asset"], box, eff_theme, it.get("recolor", False)), box[0], box[1])
        elif role == "headline":
            col = textcol(it)
            fp = _font_for(it.get("font"), _HEAD_FONT)
            if it.get("kw_title"):                       # 표지: 키워드(큰)+제목(작게, 무조건 한 줄) 원본 사이즈·정렬 유지
                kw = (headline or "").strip().rstrip(",")
                ti = " ".join((title or "").split())
                if kw and ti.startswith(kw):             # 제목이 키워드로 시작하면 중복 제거
                    ti = ti[len(kw):].lstrip(" ,·")
                if kw and not kw.endswith((",", ".", "!", "?")):
                    kw += ","
                canvas = _draw_kw_title(canvas, box, kw, ti, fp,
                                        it.get("kw_size", it.get("size", 80)), it.get("title_size", 54),
                                        col, pt2px, it.get("align", "right"),
                                        it.get("gap", 0), it.get("dy", 0))
            else:
                txt = (title or "") if it.get("source") == "title" else (headline or "")
                if it.get("lines") == 2:
                    txt = _two_lines(txt)
                canvas = _draw_text(canvas, box, txt, fp, it.get("size", 40), col, it.get("align", "center"), pt2px)
        elif role == "subtitle":
            sub = subtitle or ""
            if it.get("lines") == 2 and sub and "\n" not in sub:
                sub = _two_lines(sub)          # 서브카피 2줄 고정(길이 균형 분할)
            canvas = _draw_text(canvas, box, sub, _font_for(it.get("font"), _HEAD_FONT), it.get("size", 28), textcol(it), it.get("align", "center"), pt2px, oneline=bool(it.get("oneline")))
        elif role == "body":
            bod = body or ""
            if it.get("lines") == 2 and bod and "\n" not in bod:
                bod = _two_lines(bod)          # 내지 본문 2줄 고정(길이 균형 분할)
            canvas = _draw_text(canvas, box, bod, _font_for(it.get("font"), _BODY_FONT), it.get("size", 19), textcol(it), it.get("align", "center"), pt2px)
        elif role == "label":
            canvas = _draw_text(canvas, box, it.get("text", ""), _font_for(it.get("font"), _LABEL_FONT), it.get("size", 14), textcol(it), it.get("align", "left"), pt2px)
    canvas.convert("RGB").save(out_path, "PNG")
    return out_path


def make_cards(photo_paths, headlines, out_dir, assets_dir, theme=None, subtitle="", bodies=None, title=""):
    """카드 N장 생성. subtitle=표지부제, bodies=카드별본문, title=블로그제목(source=title 헤드라인용)."""
    os.makedirs(out_dir, exist_ok=True)
    bodies = bodies or []
    if theme is None:
        theme = PALETTE[random.choice([c for c in PALETTE if c != "orange"])]
    lay = _layout(assets_dir)
    slides = lay["slides"]
    n = len(slides)
    salt = "%08x" % random.randrange(16 ** 8)
    # 색 통일: random_theme=생성마다 랜덤색, 아니면 표지(0번) 사진 대표색
    uses_auto = lay.get("auto_theme") or any(it.get("auto_color") for sl in slides for it in sl)
    accent = None
    if lay.get("random_theme"):
        accent = _random_accent()
    elif uses_auto and photo_paths and os.path.isfile(photo_paths[0]):
        accent = _accent_from_image(photo_paths[0])
    pngs, cards = [], []
    # 헤드라인 슬라이드는 '헤드라인과 같은 인덱스의 사진'을 써서 제목-사진을 맞춘다.
    # 사진만 슬라이드(예: 퍼스트디자인 s5)는 여분 사진을 써서 뒤 카드가 밀리지 않게 한다.
    n_heads = sum(1 for sl in slides
                  if any(it.get("role") == "headline" and it.get("source") != "title" for it in sl))
    hi = 0; spare = n_heads
    for i in range(n):
        head_item = next((it for it in slides[i] if it.get("role") == "headline"), None)
        head = bod = ""
        if head_item is not None and head_item.get("source") != "title":
            head = headlines[hi] if hi < len(headlines) else ""
            bod = bodies[hi] if hi < len(bodies) else ""
            src = photo_paths[hi % len(photo_paths)] if photo_paths else ""
            hi += 1
        else:
            src = photo_paths[spare % len(photo_paths)] if photo_paths else ""
            spare += 1
        sub = subtitle if any(it.get("role") == "subtitle" for it in slides[i]) else ""
        out = os.path.join(out_dir, "card%02d.png" % (i + 1))
        render_card(i, src, head, theme, out, assets_dir, seed=f"{salt}:{i}", subtitle=sub, body=bod, title=title, accent=accent)
        flds = []
        for it in slides[i]:
            r = it.get("role")
            if r == "headline":
                flds.append("title" if it.get("source") == "title" else "headline")
            elif r in ("subtitle", "body"):
                flds.append(r)
        pngs.append(out)
        cards.append({"src": src, "headline": head, "subtitle": sub, "body": bod, "title": title,
                      "fields": flds, "cx": 0.5, "cy": 0.5, "zoom": 1.0})
    state = {"theme": list(theme), "assets_dir": assets_dir, "cards": cards,
             "accent": list(accent) if accent else None}
    json.dump(state, open(os.path.join(out_dir, "cards.json"), "w", encoding="utf-8"), ensure_ascii=False)
    return {"pngs": pngs, "theme": "%02X%02X%02X" % theme, "cards": cards}


def load_state(out_dir):
    return json.load(open(os.path.join(out_dir, "cards.json"), encoding="utf-8"))


def save_state(out_dir, state):
    json.dump(state, open(os.path.join(out_dir, "cards.json"), "w", encoding="utf-8"), ensure_ascii=False)


def edit_card(out_dir, index, src=None, center=None, zoom=None,
              headline=None, subtitle=None, body=None, title=None):
    """cards.json 상태 기반으로 카드 1장만 다시 렌더(사진교체/위치·확대/글 수정)."""
    st = load_state(out_dir); c = st["cards"][index]
    adir = st.get("assets_dir")
    if src is not None:
        c["src"] = src; c["cx"] = 0.5; c["cy"] = 0.5; c["zoom"] = 1.0
    if center is not None:
        c["cx"], c["cy"] = float(center[0]), float(center[1])
    if zoom is not None:
        c["zoom"] = float(zoom)
    if headline is not None:
        c["headline"] = headline
    if subtitle is not None:
        c["subtitle"] = subtitle
    if body is not None:
        c["body"] = body
    if title is not None:
        c["title"] = title
    out = os.path.join(out_dir, "card%02d.png" % (index + 1))
    render_card(index, c["src"], c["headline"], tuple(st["theme"]), out, adir,
                center=(c["cx"], c["cy"]), zoom=c["zoom"], seed=f"edit:{index}",
                subtitle=c.get("subtitle", ""), body=c.get("body", ""), title=c.get("title", ""),
                accent=st.get("accent"))
    save_state(out_dir, st)
    return out


def extract_template(pptx_path, out_assets_dir, theme_recolor=False, photo_round=0.0):
    """어떤 카드뉴스 PPTX든 장식 PNG + layout.json 으로 추출(브랜드 카드 등록용). 정사각 가정.
    - 사진 슬롯 = 슬라이드별 '비(非)풀블리드 이미지 중 최대 바이트'(배경 풀블리드는 고정 장식으로 유지).
    - theme_recolor=True면 하오팩토리식(주황 장식→테마색, 헤드라인=테마색). False면 디자인 그대로 + 헤드라인은 원래 색.
    - 헤드라인/라벨의 실제 글자색을 캡처."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    os.makedirs(out_assets_dir, exist_ok=True)
    prs = Presentation(pptx_path)
    W = prs.slide_width; H = prs.slide_height
    SIZE = 1080
    sc = SIZE / float(W)
    pt_height = H / 914400.0 * 72.0
    def px(v):
        return int(round(int(v) * sc)) if v is not None else 0
    def slide_bg(s):
        cSld = s._element.find(qn("p:cSld"))
        bgEl = cSld.find(qn("p:bg")) if cSld is not None else None
        if bgEl is not None:
            srgb = bgEl.findall(".//" + qn("a:srgbClr"))
            if srgb:
                return srgb[0].get("val")
        return "FFFFFF"
    spec = {"size": SIZE, "pt_height": round(pt_height, 2), "bg": "FFFFFF", "slide_bg": [],
            "theme_recolor": bool(theme_recolor), "photo_round": photo_round, "slides": []}
    for f in os.listdir(out_assets_dir):
        if f.endswith(".png") or f.endswith(".jpg") or f.endswith(".jpeg"):
            try: os.remove(os.path.join(out_assets_dir, f))
            except Exception: pass
    for i, s in enumerate(prs.slides):
        pics = []   # (shape, box, nb, part)
        texts = []
        for sh in s.shapes:
            box = [px(sh.left), px(sh.top), px(sh.width), px(sh.height)]
            if sh.shape_type == 13:
                part = s.part.related_part(sh._element.blipFill.blip.rEmbed)
                pics.append((sh, box, len(part.blob), part))
            elif sh.has_text_frame and sh.text_frame.text.strip():
                texts.append((sh, box))
        # 사진 슬롯: 비풀블리드 중 최대 바이트(없으면 전체 최대)
        def fullbleed(b):
            return b[0] <= 2 and b[1] <= 2 and b[2] >= SIZE - 5 and b[3] >= SIZE - 5
        nonfb = [p for p in pics if not fullbleed(p[1])]
        photo_shape = max(nonfb or pics, key=lambda p: p[2]) if pics else None
        if photo_shape is not None:   # 너무 작으면(타이틀카드 장식) 사진 슬롯 아님
            pb = photo_shape[1]
            if pb[2] * pb[3] < 80000:
                photo_shape = None
        items = []
        for j, (sh, box, nb, part) in enumerate(pics):
            if photo_shape is not None and sh is photo_shape[0]:
                items.append({"role": "photo", "box": box})
            else:
                ext = os.path.splitext(part.partname)[1] or ".png"
                asset = "s%d_%d%s" % (i, j, ext)
                ap = os.path.join(out_assets_dir, asset)
                with open(ap, "wb") as f:
                    f.write(part.blob)
                # 원본 그림의 투명도(alphaModFix) 반영 → 반투명 레이어 보존(불투명 통짜 방지)
                amf = sh._element.findall(".//" + qn("a:alphaModFix"))
                if asset.lower().endswith(".png") and amf and amf[0].get("amt"):
                    op = int(amf[0].get("amt")) / 100000.0
                    if op < 0.97:
                        try:
                            im = Image.open(ap).convert("RGBA")
                            r, g, b, al = im.split(); al = al.point(lambda v: int(v * op))
                            Image.merge("RGBA", (r, g, b, al)).save(ap)
                        except Exception:
                            pass
                recolor = theme_recolor and (box[3] <= 3 or not (i == 0 and box[2] < 220 and box[1] < 160))
                items.append({"role": "deco", "box": box, "asset": asset, "recolor": bool(recolor)})
        def _fsz(sh):
            r = sh.text_frame.paragraphs[0].runs
            return (r[0].font.size.pt if r and r[0].font.size else 40)
        maxsz = max((_fsz(sh) for sh, _ in texts), default=0)
        head_used = False
        for sh, box in texts:
            tf = sh.text_frame; p0 = tf.paragraphs[0]; run = p0.runs[0] if p0.runs else None
            txt = tf.text.strip()
            fsz = (run.font.size.pt if run and run.font.size else 40)
            color = None; font = None
            if run is not None:
                try:
                    if run.font.color and run.font.color.type is not None:
                        color = str(run.font.color.rgb)
                except Exception:
                    pass
                rPr = run._r.find(qn("a:rPr"))   # 한글(ea) 폰트 우선, 없으면 latin
                if rPr is not None:
                    ea = rPr.find(qn("a:ea")); la = rPr.find(qn("a:latin"))
                    font = (ea.get("typeface") if ea is not None else None) or \
                           (la.get("typeface") if la is not None else None)
                font = font or run.font.name
            # 슬라이드에서 가장 큰 글자 1개만 동적 헤드라인, 나머지는 고정 라벨(원문 유지)
            role = "headline" if (not head_used and fsz == maxsz and fsz >= 24) else "label"
            if role == "headline":
                head_used = True
            al = p0.alignment
            align = "center" if al == 2 else ("right" if al == 3 else "left")
            items.append({"role": role, "box": box, "text": (txt if role == "label" else ""),
                          "size": fsz, "color": color, "font": font, "align": align})
        spec["slide_bg"].append(slide_bg(s))
        spec["slides"].append(items)
    json.dump(spec, open(os.path.join(out_assets_dir, "layout.json"), "w", encoding="utf-8"), ensure_ascii=False, indent=1)
    _layout_cache.pop(out_assets_dir, None)
    nassets = len([f for f in os.listdir(out_assets_dir) if not f.endswith(".json")])
    return {"slides": len(spec["slides"]), "assets": nassets}
